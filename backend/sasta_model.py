import re
import AppOpener
import webbrowser
import shutil
import datetime
import random
import requests
from bs4 import BeautifulSoup
import speech_recognition as sr
import pyttsx3
from pptx import Presentation
import json
import os
import pyautogui
import time
engine = pyttsx3.init()
common_websites = [
    "google",
    "youtube",
    "facebook",
    "instagram",
    "twitter",
    "linkedin",
    "github",
    "stackoverflow",
    "amazon",
    "flipkart",
    "wikipedia",
    "netflix",
    "gmail",
    "reddit",
    "quora",
    "microsoft",
    "apple",
    "zoom",
    "whatsapp",
    "spotify",
    "chat GPT"
]
actual_website_names=[]
def automation(task):
    seperated_words = task.lower().split()
    if "open" in seperated_words:
            word_open = seperated_words.index("open")
            items_to_open = seperated_words[word_open+1:] # Get all words after "open"
            
            for item in items_to_open:
                if item in common_websites: # Check if it's a recognized item
                    if item == "whatsapp":
                        AppOpener.open(item)
                        print(f"Opening WhatsApp desktop app.")
                    else:
                        try:
                            # Try to open as a desktop app first
                            AppOpener.open(item)
                            print(f"Opening {item} desktop app.")
                        except Exception as e:
                            # If not a desktop app or fails, try opening as a website
                            print(f"Could not open {item} as a desktop app: {e}. Attempting to open as website.")
                            webbrowser.open(f"https://{item}.com")
                            print(f"Opening {item}.com in browser.")
                else:
                    print(f"'{item}' is not a recognized app or website in common_websites list.")
    else:
        print("No 'open' command detected in task.")
def general_quries(querry):
    

    jokes = [
        "Why don't scientists trust atoms? Because they make up everything!",
        "Why did the computer show up late to work? It had a hard drive!",
        "Why was the math book sad? Because it had too many problems.",
        "Parallel lines have so much in common... it's a shame they'll never meet."
    ]
    user_input = querry

    if user_input == 'bye':
        return(" Goodbye!")

    elif 'hello' in user_input or 'hi' in user_input:
        return("Hello there!")

    elif 'how are you' in user_input:
        return("I'm just a bunch of code, but I'm doing great!")

    elif 'your name' in user_input:
        return("I'm ChatBot, your friendly assistant.")

    elif 'who created you' in user_input:
        return("No one created me, haha. I just popped into existence! siuuuuuuuuuuuuuuuuuuuuuu")

    elif 'my age' in user_input:
         return("I'd guess you're 19. Just a lucky guess!")

    elif 'time' in user_input:
        current_time = datetime.datetime.now().strftime("%I:%M %p")
        return(f"The current time is {current_time}.")

    elif 'date' in user_input or 'today\'s date' in user_input:
        current_date = datetime.datetime.now().strftime("%B %d, %Y")
        return(f"Today's date is {current_date}.")

    elif 'day of the week' in user_input or 'what day is it' in user_input:
        day_name = datetime.datetime.now().strftime("%A")
        return(f"Today is {day_name}.")

    elif 'joke' in user_input:
        return(f" {random.choice(jokes)}")

    elif 'weather' in user_input:
        return("It's always sunny in code world!")

    elif 'thank you' in user_input or 'thanks' in user_input:
        return("You're welcome!")

    elif 'favourite color' in user_input or 'favorite color' in user_input:
        return("I like blue – reminds me of the terminal screen!")

    elif 'favourite food' in user_input or 'favorite food' in user_input:
        return("Binary bites and code crunchers 😄")

    else:
        return("ChatBot: Sorry, I didn't understand that.")
GROQ_API_KEY = "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"
GROQ_ENDPOINT = "https://api.groq.com/openai/v1/chat/completions"
## clean responsse code is of gpt bcs the ask_groq functionn was returning random shit so i wanted to clean the response
def clean_response(text):
    # Remove code blocks, markdown symbols, extra spaces, etc.
    text = re.sub(r"```.*?```", "", text, flags=re.DOTALL)  # remove code blocks
    text = re.sub(r"[*_`#>-]", "", text)  # remove markdown formatting
    text = re.sub(r"\s+", " ", text)  # collapse whitespace
    text = re.sub(r"[^\x00-\x7F]+", "", text)  # remove non-ASCII characters (if TTS struggles)
    text = text.strip()
    return text
def brave_search(query):
    headers = {
        "User-Agent": "Mozilla/5.0"
    }
    search_url = f"https://search.brave.com/search?q={query.replace(' ', '+')}"
    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    results = []
    for result in soup.find_all("div", class_="snippet"):
        snippet_text = result.get_text()
        results.append(snippet_text)

    return results[:3]  
def ask_groq(question):
    headers = {
    "Authorization": f"Bearer {GROQ_API_KEY}",
    "Content-Type": "application/json"
}

    json_data = {
    "model": "llama-3.3-70b-versatile",
    "messages": [
        {
            "role": "system",
            "content": (
                f"answer only to this  {question} dont give unnecessay answers"
            )
        },
        {
            "role": "user",
            "content": f"{question}"
        }
    ]
}
    response = requests.post(GROQ_ENDPOINT, headers=headers, json=json_data)
    response.raise_for_status()
    raw_output = response.json()['choices'][0]['message']['content']
    return clean_response(raw_output)
def listen_and_convert():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)
        try:
            text = recognizer.recognize_google(audio)
            print(" You said:", text)
            return text
        except sr.UnknownValueError:
            return(" Could not understand audio")
        except sr.RequestError as e:
            print("Could not request results; {0}".format(e))
def txt_to_speak(queery):
    engine.say(queery)
    engine.runAndWait()
def generate_content_for_ppt(prompt):
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }
    json_data = {
        "model": "llama-3.3-70b-versatile",
        "messages": [
            {
                "role": "system",
                "content": (
                    f"Generate high-quality and meaningful slide titles, contents, and image URLs for the topic: {prompt}. "
                    "Return a valid JSON object with keys: 'titles', 'contents', and 'images'. No explanations, just JSON."
                )
            },
            {
                "role": "user",
                "content": prompt
            }
        ]
    }
    response = requests.post(GROQ_ENDPOINT, headers=headers, json=json_data)
    response.raise_for_status()
    raw_response = response.json()['choices'][0]['message']['content'].strip()
    if raw_response.startswith("```json"):
        raw_response = raw_response[7:-3].strip()
    try:
        actual_response = json.loads(raw_response)
        titles = actual_response.get("titles", [])
        contents = actual_response.get("contents", [])
        images = actual_response.get("images", [])
        return titles, contents, images
    except json.JSONDecodeError as e:
        print("Failed to decode JSON:", e)
        return [], [], []
def sasta_ppt_maker(ppt_name, titles, contents):
    print("content",contents,titles)
    prs = Presentation()
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide0.shapes.title
    subtitle = slide0.placeholders[1]
    title.text = "Welcome to Jarvis"
    subtitle.text = "sasta_assistant"    
    for i in range(len(titles)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = slide.shapes.title
        slide_content = slide.placeholders[1]
        slide_title.text = titles[i]
        slide_content.text = contents[i]
    save_path = os.path.join(r"C:\Users\PRERAN S\OneDrive\Desktop", ppt_name)
    prs.save(save_path)
    print(f"PPT saved at: {save_path}")
def generate_code(question):
    headers = {
    "Authorization": f"Bearer {GROQ_API_KEY}",
    "Content-Type": "application/json"
}

    json_data = {
    "model": "llama-3.3-70b-versatile",
    "messages": [
        {
            "role": "system",
            "content": (
                f"generate a apporiate code for the{question} in very presice format add comments to understand the code in python lanuage"
            )
        },
        {
            "role": "user",
            "content": f"{question}"
        }
    ]
}
    response = requests.post(GROQ_ENDPOINT, headers=headers, json=json_data)
    response.raise_for_status()
    result = response.json()
    return result['choices'][0]['message']['content']

def whatsapp_sasta_automation(to,message):
    AppOpener.open("whatsapp")
    time.sleep(5)
    pyautogui.click(x=200, y=150)  
    time.sleep(1)
    pyautogui.write(f"{to}", interval=0.1)
    time.sleep(1)
    pyautogui.click(x=200, y=200)  
    time.sleep(1)
    pyautogui.write(f"{message}", interval=0.1)
    pyautogui.press("enter")
def youtube_sasta_automation(channel):
    webbrowser.open("https://youtube.com")
    time.sleep(5)
    pyautogui.click(x=950, y=150)  
    time.sleep(1)
    pyautogui.write(channel,interval=0.1)
    time.sleep(1)
    pyautogui.press("enter")
    time.sleep(1)
    pyautogui.click(x=600, y=550)  
#######################################################################################################################################################
text =listen_and_convert()

general_pattern = r"\b(hello|hi|how|are|you|thanks|goodbye|bye|introduce|joke|name|date|time|fact|space|explain|describe|define|mean|meaning|info|information|question|ask)\b"
automation_pattern = r"\b(open|launch|start|begin|initiate|close|shut|turn|on|off|play|stop|pause|resume|call|dial|message|text|email|remind|set|schedule|create|make|run|search|navigate|go|find|install|uninstall|delete|send|post)\b"
realtime_pattern = r"\b(weather|temperature|forecast|humidity|air|quality|pollution|price|rate|value|stock|market|today|now|current|news|headlines|traffic|covid|cases|score|match|matches|exchange|convert|currency|bitcoin|crypto|population|sunrise|sunset|trending|live|update|event|events|happening|ongoing|stats|statistics|standing|rank|ranking|winner|result|results|ipl|election|vote|alert|signal|data|how|what|who|when|why|where)\b"

ppt_format = r"generate\s+a\s+ppt\s+on\s+topic\s+(.+?)\s+and\s+save\s+it\s+as\s+(\w+)"
whatsapp_pattern = r"open whatsapp\s+and\s+send\s+message\s+to\s+(\w+)\s+saying\s+(.+)"
youtube_pattern = r"open youtube\s+and\s+search\s+for\s+(.+)"
sasta_search_optimization = r"\b(open|launch|start)\s+(\w+)(?:\s+(?:and\s+)?search\s+for\s+(.+))?"
executed = False

if match := re.search(ppt_format, text, re.IGNORECASE):
    topic, filename = match.groups()
    filename += ".pptx"
    titles, contents, images = generate_content_for_ppt(topic)
    sasta_ppt_maker(filename, titles, contents)
    executed = True
elif match := re.search(whatsapp_pattern, text, re.IGNORECASE):
    whatsapp_sasta_automation(*match.groups())
    executed = True
elif match := re.search(youtube_pattern, text, re.IGNORECASE):
    youtube_sasta_automation(match.group(1))
    executed = True
elif match := re.search(sasta_search_optimization, text, re.IGNORECASE):
    app, _, query = match.groups()
    if app and query:
        webbrowser.open(f"https://www.{app}.com/results?search_query={query}")
        executed = True
if not executed:
    chunks = re.split(r"\b(?:and then|then|also|finally|after that|and)\b", text, flags=re.IGNORECASE)
    chunks = [chunk.strip() for chunk in chunks if chunk.strip()]
    general, automation_tasks, realtime = [], [], []
    for chunk in chunks:
        g = len(re.findall(general_pattern, chunk))
        a = len(re.findall(automation_pattern, chunk))
        r = len(re.findall(realtime_pattern, chunk))
        if r >= g and r >= a:
            realtime.append(chunk)
        elif a >= g and a >= r:
            automation_tasks.append(chunk)
        else:
            general.append(chunk)
    for task in general:
        result = general_quries(task)
        txt_to_speak(result)
    for task in automation_tasks:
        automation(task)
    for task in realtime:
        search_result = brave_search(task)
        final_answer = ask_groq(search_result)
        txt_to_speak(final_answer)
