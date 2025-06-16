import os
from pptx import Presentation
import requests
import json
GROQ_API_KEY = "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"
GROQ_ENDPOINT = "https://api.groq.com/openai/v1/chat/completions"
def generate_content_for_ppt(prompt):
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }
    json_data = {
        "model": "llama-3.3-70b-versatile",
        "messages": [
            {
                "role": "system",
                "content": (
                    f"Generate high-quality and meaningful slide titles, contents, and image URLs for the topic: {prompt}. "
                    "Return a valid JSON object with keys: 'titles', 'contents', and 'images'. No explanations, just JSON."
                )
            },
            {
                "role": "user",
                "content": prompt
            }
        ]
    }
    response = requests.post(GROQ_ENDPOINT, headers=headers, json=json_data)
    response.raise_for_status()
    raw_response = response.json()['choices'][0]['message']['content'].strip()
    if raw_response.startswith("```json"):
        raw_response = raw_response[7:-3].strip()
    try:
        actual_response = json.loads(raw_response)
        titles = actual_response.get("titles", [])
        contents = actual_response.get("contents", [])
        images = actual_response.get("images", [])
        return titles, contents, images
    except json.JSONDecodeError as e:
        print("Failed to decode JSON:", e)
        return [], [], []
def sasta_ppt_maker(ppt_name, titles, contents):
    print("content",contents,titles)
    prs = Presentation()
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide0.shapes.title
    subtitle = slide0.placeholders[1]
    title.text = "Welcome to Jarvis"
    subtitle.text = "sasta_assistant"    
    for i in range(len(titles)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = slide.shapes.title
        slide_content = slide.placeholders[1]
        slide_title.text = titles[i]
        slide_content.text = contents[i]
    save_path = os.path.join(r"C:\Users\PRERAN S\OneDrive\Desktop", ppt_name)
    prs.save(save_path)
    print(f"PPT saved at: {save_path}")

